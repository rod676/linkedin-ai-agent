from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import json
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Design System — Couleurs
COLORS = {
    "background_dark": RGBColor(10, 10, 20),
    "background_card": RGBColor(20, 25, 40),
    "accent_blue": RGBColor(0, 120, 255),
    "accent_cyan": RGBColor(0, 200, 200),
    "text_white": RGBColor(255, 255, 255),
    "text_gray": RGBColor(160, 170, 190),
    "danger_red": RGBColor(220, 50, 50),
    "success_green": RGBColor(50, 200, 100),
}

# Format LinkedIn carrousel — carré
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(10)


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    return RGBColor(r, g, b)


def add_background(slide, prs, color: RGBColor):
    """Ajoute un fond coloré à la slide."""
    from pptx.util import Inches
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    shape = slide.shapes.add_shape(
        1, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=24, bold=False,
                 color=None, align=PP_ALIGN.LEFT):
    """Ajoute une text box stylée."""
    if color is None:
        color = COLORS["text_white"]
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    
    return txBox


def add_accent_line(slide, prs):
    """Ajoute une ligne décorative bleue."""
    from pptx.util import Inches, Pt
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(1.5)
    height = Inches(0.05)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["accent_blue"]
    shape.line.fill.background()


def create_hook_slide(slide, prs, slide_data: dict, total_slides: int):
    """Slide 1 — Hook accrocheur."""
    add_background(slide, prs, COLORS["background_dark"])
    add_accent_line(slide, prs)
    
    # Emoji grand
    add_text_box(
        slide,
        slide_data.get("emoji", "🔐"),
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(1.5),
        font_size=60, align=PP_ALIGN.CENTER
    )
    
    # Titre principal
    add_text_box(
        slide,
        slide_data["headline"],
        Inches(0.5), Inches(3.2),
        Inches(9), Inches(3),
        font_size=40, bold=True,
        color=COLORS["text_white"],
        align=PP_ALIGN.CENTER
    )
    
    # Sous-titre
    if slide_data.get("body"):
        add_text_box(
            slide,
            slide_data["body"],
            Inches(0.5), Inches(6.5),
            Inches(9), Inches(1.5),
            font_size=22,
            color=COLORS["text_gray"],
            align=PP_ALIGN.CENTER
        )
    
    # Swipe indicator
    add_text_box(
        slide,
        "Swipe to learn more →",
        Inches(0.5), Inches(8.5),
        Inches(9), Inches(0.8),
        font_size=18,
        color=COLORS["accent_cyan"],
        align=PP_ALIGN.CENTER
    )
    
    # Numéro slide
    add_text_box(
        slide,
        f"1/{total_slides}",
        Inches(8.5), Inches(0.3),
        Inches(1), Inches(0.5),
        font_size=14,
        color=COLORS["text_gray"],
        align=PP_ALIGN.RIGHT
    )


def create_insight_slide(slide, prs, slide_data: dict,
                          slide_num: int, total_slides: int):
    """Slides 2-N — Insights."""
    add_background(slide, prs, COLORS["background_dark"])
    
    # Barre latérale bleue
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        Inches(0.15), Inches(10)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["accent_blue"]
    bar.line.fill.background()
    
    # Emoji
    add_text_box(
        slide,
        slide_data.get("emoji", "⚡"),
        Inches(0.5), Inches(1),
        Inches(9), Inches(1.2),
        font_size=48, align=PP_ALIGN.LEFT
    )
    
    # Headline
    add_text_box(
        slide,
        slide_data["headline"],
        Inches(0.5), Inches(2.5),
        Inches(9), Inches(2.5),
        font_size=34, bold=True,
        color=COLORS["text_white"]
    )
    
    # Body
    add_text_box(
        slide,
        slide_data.get("body", ""),
        Inches(0.5), Inches(5.2),
        Inches(9), Inches(3),
        font_size=22,
        color=COLORS["text_gray"]
    )
    
    # Numéro
    add_text_box(
        slide,
        f"{slide_num}/{total_slides}",
        Inches(8.5), Inches(0.3),
        Inches(1), Inches(0.5),
        font_size=14,
        color=COLORS["text_gray"],
        align=PP_ALIGN.RIGHT
    )


def create_cta_slide(slide, prs, carousel_data: dict, total_slides: int):
    """Dernière slide — CTA."""
    add_background(slide, prs, COLORS["accent_blue"])
    
    # Titre CTA
    add_text_box(
        slide,
        "Found this useful?",
        Inches(0.5), Inches(2),
        Inches(9), Inches(1.5),
        font_size=38, bold=True,
        color=COLORS["text_white"],
        align=PP_ALIGN.CENTER
    )
    
    # Actions
    add_text_box(
        slide,
        "♻️ Repost to help your network\n❤️ Like if you learned something\n💬 Comment your thoughts below",
        Inches(0.5), Inches(3.8),
        Inches(9), Inches(3),
        font_size=24,
        color=COLORS["text_white"],
        align=PP_ALIGN.CENTER
    )
    
    # Author
    add_text_box(
        slide,
        carousel_data.get("author", ""),
        Inches(0.5), Inches(7.5),
        Inches(9), Inches(1),
        font_size=18,
        color=COLORS["text_white"],
        align=PP_ALIGN.CENTER
    )
    
    # Numéro
    add_text_box(
        slide,
        f"{total_slides}/{total_slides}",
        Inches(8.5), Inches(0.3),
        Inches(1), Inches(0.5),
        font_size=14,
        color=COLORS["text_white"],
        align=PP_ALIGN.RIGHT
    )


def generate_carousel_pptx(carousel_data: dict,
                             output_path: str = "posts/carousel.pptx") -> str:
    """Génère le fichier PPTX du carrousel."""
    
    os.makedirs("posts", exist_ok=True)
    
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    blank_layout = prs.slide_layouts[6]
    slides_data = carousel_data["slides"]
    total_slides = len(slides_data) + 1  # +1 pour CTA
    
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        
        if slide_data["type"] == "hook":
            create_hook_slide(slide, prs, slide_data, total_slides)
        else:
            create_insight_slide(
                slide, prs, slide_data,
                i + 1, total_slides
            )
    
    # CTA slide
    cta_slide = prs.slides.add_slide(blank_layout)
    create_cta_slide(cta_slide, prs, carousel_data, total_slides)
    
    prs.save(output_path)
    print(f"✅ Carousel saved: {output_path}")
    return output_path


def generate_carousel_content(topic: str) -> dict:
    """GPT génère le contenu structuré du carrousel."""
    from prompts import carousel_prompt
    from researcher import research_topic
    
    print(f"🔍 Researching: {topic}")
    facts = research_topic(topic)
    
    print(f"🧠 Generating carousel structure...")
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {
                "role": "system",
                "content": "You are an expert LinkedIn carousel creator for Cybersecurity & AI content. Always respond in English."
            },
            {
                "role": "user",
                "content": carousel_prompt(topic, facts)
            }
        ],
        temperature=0.7,
        max_tokens=1500
    )
    
    raw = response.choices[0].message.content.strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    
    return json.loads(raw)